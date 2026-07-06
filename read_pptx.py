import sys
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*20} Slide {i+1} {'='*20}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text.strip())
                print("-" * 10)

if __name__ == "__main__":
    extract_text(r"C:\Users\dvsat\Desktop\MSL\webDesign.pptx")
